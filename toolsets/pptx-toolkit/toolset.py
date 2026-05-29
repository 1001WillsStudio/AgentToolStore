"""
pptx‑toolkit — Create, edit, and extract content from PowerPoint (.pptx) presentations.
============================================================================================

Slide‑level operations: read text, inspect structure, extract speaker notes,
and create presentations from structured content.
"""

from pathlib import Path

try:
    from toolstore.toolset import tool
except ImportError:
    def tool(fn):
        return fn


# ── pptx_read ──────────────────────────────────────────────────────────

@tool
def pptx_read(*, filepath: str) -> dict:
    """Extract all text from a PowerPoint presentation, slide by slide.

    Args:
        filepath: Path to the .pptx file.

    Returns:
        dict with:
          slides — list of {index, title, body_text, notes, shape_count}
          count  — total number of slides
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx not installed — run: pip install python-pptx"}

    p = Path(filepath).expanduser().resolve()
    if not p.exists():
        return {"error": f"File not found: {p}"}

    try:
        prs = Presentation(str(p))
    except Exception as exc:
        return {"error": f"Cannot open presentation: {exc}"}

    slides = []
    for i, slide in enumerate(prs.slides):
        title = ""
        body_parts = []
        notes_text = ""

        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                title = shape.text_frame.text if shape.has_text_frame else ""
            elif shape.has_text_frame:
                body_parts.append(shape.text_frame.text)
            elif shape.has_table:
                tbl_text = []
                for row in shape.table.rows:
                    tbl_text.append(" | ".join(cell.text for cell in row.cells))
                body_parts.append("\n".join(tbl_text))

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text

        slides.append({
            "index": i + 1,
            "title": title.strip(),
            "body_text": "\n".join(bp for bp in body_parts if bp.strip()),
            "notes": notes_text.strip(),
            "shape_count": len(slide.shapes),
        })

    return {"slides": slides, "count": len(slides)}


# ── pptx_info ──────────────────────────────────────────────────────────

@tool
def pptx_info(*, filepath: str) -> dict:
    """Inspect presentation structure: slide count, layout, dimensions.

    Args:
        filepath: Path to the .pptx file.

    Returns:
        dict with: slides, slide_width, slide_height, layouts, filename
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx not installed — run: pip install python-pptx"}

    p = Path(filepath).expanduser().resolve()
    if not p.exists():
        return {"error": f"File not found: {p}"}

    try:
        prs = Presentation(str(p))
    except Exception as exc:
        return {"error": f"Cannot open presentation: {exc}"}

    layouts = [ly.name for ly in prs.slide_layouts]

    return {
        "slides": len(prs.slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "layouts": layouts,
        "filename": p.name,
    }


# ── pptx_create ────────────────────────────────────────────────────────

@tool
def pptx_create(*, filepath: str, title: str = "",
                slides_data: list = None) -> dict:
    """Create a new PowerPoint presentation from structured slide data.

    Args:
        filepath:    Where to write the new .pptx file.
        title:       Presentation title (used as first slide).
        slides_data: List of dicts: [{title, bullets: [...]}, ...].

    Returns:
        dict with "written" (absolute path) or "error".
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return {"error": "python-pptx not installed — run: pip install python-pptx"}

    p = Path(filepath).expanduser().resolve()
    p.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    # Title slide
    if title:
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = title

    # Content slides
    for sd in (slides_data or []):
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title and sd.get("title"):
            slide.shapes.title.text = sd["title"]

        # Add bullets
        bullets = sd.get("bullets", [])
        if bullets and len(slide.shapes) > 1:
            body_shape = slide.shapes[1]
            if body_shape.has_text_frame:
                tf = body_shape.text_frame
                tf.clear()
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        tf.text = bullet
                    else:
                        p2 = tf.add_paragraph()
                        p2.text = bullet
                        p2.level = 0

    prs.save(str(p))
    return {"written": str(p), "slides": len(prs.slides)}
